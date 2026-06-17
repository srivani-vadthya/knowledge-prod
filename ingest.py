from openai import OpenAI
from dotenv import load_dotenv
from pinecone_client import get_index
from pypdf import PdfReader
from pypdf.errors import PdfReadError
import fitz
import os
from docx import Document
import csv
import openpyxl
import pandas as pd
from pptx import Presentation
import markdown
from bs4 import BeautifulSoup
from PIL import Image
import io
import base64
try:
    import easyocr
    EASYOCR_AVAILABLE = True
except ImportError:
    EASYOCR_AVAILABLE = False

try:
    import pytesseract
    PYTESSERACT_AVAILABLE = True
except ImportError:
    PYTESSERACT_AVAILABLE = False

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
load_dotenv(os.path.join(BASE_DIR, ".env"))

CHUNK_SIZE = 800
CHUNK_OVERLAP = 100
EMBED_BATCH_SIZE = 64
UPSERT_BATCH_SIZE = 100
OCR_ENGINE = os.getenv("OCR_ENGINE", "easyocr")  # easyocr or tesseract
VISION_MODEL = os.getenv("VISION_MODEL", "gpt-4o-mini")  # OpenAI vision model
ENABLE_IMAGE_PROCESSING = os.getenv("ENABLE_IMAGE_PROCESSING", "true").lower() == "true"


def extract_images_from_pdf(file_path):
    """Extract images from PDF pages"""
    images = []
    try:
        pdf_document = fitz.open(file_path)
        for page_num in range(len(pdf_document)):
            page = pdf_document[page_num]
            image_list = page.get_images(full=True)
            
            for img_index, img in enumerate(image_list):
                xref = img[0]
                base_image = pdf_document.extract_image(xref)
                image_bytes = base_image["image"]
                
                # Convert to PIL Image
                image = Image.open(io.BytesIO(image_bytes))
                
                # Skip very small images (likely icons/logos)
                if image.width < 100 or image.height < 100:
                    continue
                
                images.append({
                    "image": image,
                    "page": page_num + 1,
                    "index": img_index,
                })
        
        pdf_document.close()
    except Exception as e:
        print(f"Error extracting images: {str(e)}")
    
    return images


def ocr_image(image):
    """Extract text from image using OCR"""
    text = ""
    
    try:
        if OCR_ENGINE == "easyocr" and EASYOCR_AVAILABLE:
            reader = easyocr.Reader(['en'], gpu=False)
            # Convert PIL Image to numpy array
            import numpy as np
            img_array = np.array(image)
            results = reader.readtext(img_array, detail=0)
            text = " ".join(results)
        
        elif OCR_ENGINE == "tesseract" and PYTESSERACT_AVAILABLE:
            text = pytesseract.image_to_string(image)
        
        else:
            # Fallback: no OCR available
            text = ""
    
    except Exception as e:
        print(f"OCR error: {str(e)}")
        text = ""
    
    return text.strip()


def describe_image_with_vision(image, client):
    """Use OpenAI Vision API to describe image content"""
    try:
        # Convert PIL Image to base64
        buffered = io.BytesIO()
        image.save(buffered, format="PNG")
        img_base64 = base64.b64encode(buffered.getvalue()).decode('utf-8')
        
        response = client.chat.completions.create(
            model=VISION_MODEL,
            messages=[
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": "Describe this image in detail. Focus on any text, diagrams, charts, tables, or important visual information that would be useful for document understanding."
                        },
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/png;base64,{img_base64}"
                            }
                        }
                    ]
                }
            ],
            max_tokens=500
        )
        
        return response.choices[0].message.content.strip()
    
    except Exception as e:
        print(f"Vision API error: {str(e)}")
        return ""


def process_pdf_with_images(file_path, client):
    """Process PDF extracting both text and image content"""
    documents = []
    
    print(f"[IMAGE PROCESSING] Processing PDF: {os.path.basename(file_path)}")
    
    # Extract regular text
    try:
        pdf = fitz.open(file_path)
        text_pages = 0
        for page_number, page in enumerate(pdf, start=1):
            text = page.get_text("text").strip()
            if text:
                documents.append({
                    "text": text,
                    "metadata": {
                        "source": os.path.basename(file_path),
                        "page": page_number,
                        "type": "text"
                    },
                })
                text_pages += 1
        pdf.close()
        print(f"[IMAGE PROCESSING] Extracted text from {text_pages} pages")
    except Exception as e:
        print(f"Error extracting text: {str(e)}")
    
    # Extract and process images
    if ENABLE_IMAGE_PROCESSING:
        print(f"[IMAGE PROCESSING] Extracting images...")
        images = extract_images_from_pdf(file_path)
        print(f"[IMAGE PROCESSING] Found {len(images)} images (>100x100px)")
        
        for i, img_data in enumerate(images, 1):
            image = img_data["image"]
            page_num = img_data["page"]
            
            print(f"[IMAGE PROCESSING] Processing image {i}/{len(images)} - Page {page_num}, Size: {image.width}x{image.height}")
            
            # Run OCR
            print(f"[IMAGE PROCESSING]   → Running OCR ({OCR_ENGINE})...")
            ocr_text = ocr_image(image)
            if ocr_text:
                print(f"[IMAGE PROCESSING]   → OCR extracted {len(ocr_text)} characters")
            else:
                print(f"[IMAGE PROCESSING]   → OCR: No text detected")
            
            # Run Vision model
            print(f"[IMAGE PROCESSING]   → Running Vision API ({VISION_MODEL})...")
            vision_description = describe_image_with_vision(image, client)
            if vision_description:
                print(f"[IMAGE PROCESSING]   → Vision API returned {len(vision_description)} characters")
            else:
                print(f"[IMAGE PROCESSING]   → Vision API: No description")
            
            # Combine OCR and Vision results
            combined_text = ""
            if ocr_text:
                combined_text += f"OCR Text: {ocr_text}\n\n"
            if vision_description:
                combined_text += f"Image Description: {vision_description}"
            
            if combined_text.strip():
                documents.append({
                    "text": combined_text,
                    "metadata": {
                        "source": os.path.basename(file_path),
                        "page": page_num,
                        "type": "image",
                        "image_index": img_data["index"]
                    },
                })
                print(f"[IMAGE PROCESSING]   ✓ Image content indexed")
    
    print(f"[IMAGE PROCESSING] Complete: {len(documents)} total documents (text + images)")
    return documents


def load_documents(file_path):
    """Load documents from various file formats"""
    
    # PDF Processing with image support
    if file_path.endswith(".pdf"):
        # Use enhanced image processing if enabled
        if ENABLE_IMAGE_PROCESSING:
            client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))
            documents = process_pdf_with_images(file_path, client)
            if documents:
                return documents
        
        # Fallback to standard PDF processing
        documents = []
        try:
            pdf = fitz.open(file_path)
            for page_number, page in enumerate(pdf, start=1):
                text = page.get_text("text").strip()
                if text:
                    documents.append({
                        "text": text,
                        "metadata": {
                            "source": os.path.basename(file_path),
                            "page": page_number,
                        },
                    })
            pdf.close()
        except Exception:
            documents = []

        if documents:
            return documents

        try:
            reader = PdfReader(file_path, strict=False)
        except PdfReadError as exc:
            raise ValueError(
                "This PDF appears to be malformed and could not be read. "
                "Try opening it in a PDF viewer and printing/saving it as a new PDF, "
                "or upload a text version of the document."
            ) from exc

        for page_number, page in enumerate(reader.pages, start=1):
            try:
                text = page.extract_text() or ""
            except Exception:
                text = ""

            text = text.strip()
            if text:
                documents.append({
                    "text": text,
                    "metadata": {
                        "source": os.path.basename(file_path),
                        "page": page_number,
                    },
                })

        if not documents:
            raise ValueError(
                "No readable text was found in this PDF. If it is scanned or image-only, "
                "convert it with OCR first or upload a TXT file."
            )

        return documents
    
    # Word Document Processing
    elif file_path.endswith(".docx"):
        doc = Document(file_path)
        text = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
        if not text.strip():
            raise ValueError("No readable text found in the Word document.")
        return [{"text": text, "metadata": {"source": os.path.basename(file_path)}}]
    
    # Excel Processing (.xlsx, .xls)
    elif file_path.endswith((".xlsx", ".xls")):
        try:
            # Read all sheets
            excel_file = pd.ExcelFile(file_path)
            all_text = []
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                
                # Convert DataFrame to readable text
                sheet_text = f"Sheet: {sheet_name}\n"
                sheet_text += df.to_string(index=False)
                all_text.append(sheet_text)
            
            combined_text = "\n\n".join(all_text)
            
            if not combined_text.strip():
                raise ValueError("No readable data found in the Excel file.")
            
            return [{
                "text": combined_text,
                "metadata": {
                    "source": os.path.basename(file_path),
                    "sheets": len(excel_file.sheet_names)
                }
            }]
        except Exception as e:
            raise ValueError(f"Error reading Excel file: {str(e)}")
    
    # PowerPoint Processing (.pptx, .ppt)
    elif file_path.endswith((".pptx", ".ppt")):
        try:
            prs = Presentation(file_path)
            documents = []
            
            for slide_number, slide in enumerate(prs.slides, start=1):
                slide_text = []
                
                # Extract text from all shapes in the slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                if slide_text:
                    documents.append({
                        "text": "\n".join(slide_text),
                        "metadata": {
                            "source": os.path.basename(file_path),
                            "slide": slide_number,
                        },
                    })
            
            if not documents:
                raise ValueError("No readable text found in the PowerPoint file.")
            
            return documents
        except Exception as e:
            raise ValueError(f"Error reading PowerPoint file: {str(e)}")
    
    # Markdown Processing (.md)
    elif file_path.endswith(".md"):
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                md_text = f.read()
            
            # Convert Markdown to HTML, then extract plain text
            html = markdown.markdown(md_text)
            soup = BeautifulSoup(html, 'html.parser')
            text = soup.get_text()
            
            if not text.strip():
                raise ValueError("No readable text found in the Markdown file.")
            
            return [{
                "text": text,
                "metadata": {"source": os.path.basename(file_path)}
            }]
        except Exception as e:
            raise ValueError(f"Error reading Markdown file: {str(e)}")
    
    # CSV Processing
    elif file_path.endswith(".csv"):
        rows = []
        with open(file_path, "r", encoding="utf-8") as f:
            reader = csv.reader(f)
            rows = list(reader)
        if not rows:
            raise ValueError("CSV file is empty.")
        text = "\n".join([", ".join(row) for row in rows])
        return [{"text": text, "metadata": {"source": os.path.basename(file_path)}}]
    
    # Plain Text Processing (.txt)
    else:
        with open(file_path, "r", encoding="utf-8") as f:
            return [{
                "text": f.read(),
                "metadata": {"source": os.path.basename(file_path)},
            }]


def split_documents(documents):
    chunks = []

    for document in documents:
        text = document["text"]
        start = 0
        while start < len(text):
            end = min(start + CHUNK_SIZE, len(text))
            chunk_text = text[start:end].strip()
            if chunk_text:
                chunks.append({
                    "text": chunk_text,
                    "metadata": document["metadata"],
                })

            if end == len(text):
                break
            start = max(end - CHUNK_OVERLAP, start + 1)

    return chunks

def batch_items(items, batch_size):
    for start in range(0, len(items), batch_size):
        yield items[start:start + batch_size]

def ingest_document(file_path):
    docs = load_documents(file_path)
    chunks = split_documents(docs)
    if not chunks:
        raise ValueError("No text chunks were created from this document.")

    client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

    index = get_index()

    vectors = []
    for batch_start, chunk_batch in enumerate(batch_items(chunks, EMBED_BATCH_SIZE)):
        response = client.embeddings.create(
            input=[chunk["text"] for chunk in chunk_batch],
            model="text-embedding-3-small"
        )
        for offset, embedding_item in enumerate(response.data):
            chunk_index = batch_start * EMBED_BATCH_SIZE + offset
            chunk = chunk_batch[offset]
            vectors.append({
                "id": f"{os.path.basename(file_path)}_{chunk_index}",
                "values": embedding_item.embedding,
                "metadata": {
                    "text": chunk["text"],
                    **chunk["metadata"],
                },
            })

    for vector_batch in batch_items(vectors, UPSERT_BATCH_SIZE):
        index.upsert(vectors=vector_batch)

    return len(vectors)
